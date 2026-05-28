from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

DETAIL_OUTPUT_NAME = "education-system-framework-project-overview-2026-05-27.pptx"
EXECUTIVE_RU_OUTPUT_NAME = "education-system-framework-executive-ru-2026-05-27.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(14, 42, 71)
BLUE = RGBColor(26, 95, 180)
TEAL = RGBColor(0, 128, 128)
LIGHT = RGBColor(243, 247, 251)
MID = RGBColor(96, 110, 123)
DARK = RGBColor(32, 41, 52)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(46, 125, 50)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(198, 40, 40)

DETAIL_FOOTER = "Education System Framework • Project overview • Snapshot: 2026-05-27"
EXECUTIVE_RU_FOOTER = "Education Framework System • Executive deck (RU) • Snapshot: 2026-05-27"


def set_run_style(run, size=20, bold=False, color=DARK):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


class DeckBuilder:
    def __init__(self, footer_text):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.footer_text = footer_text

    def blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = WHITE
        return slide

    def add_header(self, slide, title, subtitle=None):
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.95))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.color.rgb = NAVY

        title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(10.8), Inches(0.34))
        p = title_box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_run_style(run, size=24, bold=True, color=WHITE)

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.47), Inches(0.54), Inches(10.6), Inches(0.2))
            p2 = subtitle_box.text_frame.paragraphs[0]
            run2 = p2.add_run()
            run2.text = subtitle
            set_run_style(run2, size=10, color=RGBColor(220, 230, 240))

        self.add_footer(slide)

    def add_footer(self, slide):
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.35), Inches(7.08), Inches(12.6), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(220, 226, 232)
        line.line.color.rgb = RGBColor(220, 226, 232)
        footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.12), Inches(12.0), Inches(0.2))
        p = footer.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = self.footer_text
        set_run_style(run, size=9, color=MID)

    def add_title_slide(self, title, subtitle_lines, callout_text):
        slide = self.blank_slide()
        banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.8), Inches(12.1), Inches(5.6))
        banner.fill.solid()
        banner.fill.fore_color.rgb = LIGHT
        banner.line.color.rgb = RGBColor(220, 226, 232)

        accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.8), Inches(12.1), Inches(0.32))
        accent.fill.solid()
        accent.fill.fore_color.rgb = BLUE
        accent.line.color.rgb = BLUE

        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10.8), Inches(1.3))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        set_run_style(run, size=28, bold=True, color=NAVY)

        sub = slide.shapes.add_textbox(Inches(1.0), Inches(2.85), Inches(10.6), Inches(2.0))
        tf2 = sub.text_frame
        tf2.word_wrap = True
        for idx, line in enumerate(subtitle_lines):
            p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
            run = p.add_run()
            run.text = line
            set_run_style(run, size=18 if idx == 0 else 16, color=DARK)
            p.space_after = Pt(10)

        callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(5.8), Inches(0.7))
        callout.fill.solid()
        callout.fill.fore_color.rgb = TEAL
        callout.line.color.rgb = TEAL
        tf3 = callout.text_frame
        tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf3.paragraphs[0].add_run()
        r.text = callout_text
        set_run_style(r, size=14, bold=True, color=WHITE)

        self.add_footer(slide)

    def add_bullet_slide(self, title, bullets, subtitle=None):
        slide = self.blank_slide()
        self.add_header(slide, title, subtitle)
        body = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(11.9), Inches(5.55))
        tf = body.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for idx, item in enumerate(bullets):
            level = 0
            text = item
            if isinstance(item, tuple):
                level, text = item
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.level = level
            p.space_after = Pt(8)
            p.bullet = True
            run = p.add_run()
            run.text = text
            set_run_style(run, size=20 if level == 0 else 16, bold=(level == 0 and idx == 0), color=DARK)

    def add_two_column_slide(self, title, left_title, left_bullets, right_title, right_bullets, subtitle=None):
        slide = self.blank_slide()
        self.add_header(slide, title, subtitle)

        for x, section_title, bullets in [
            (0.7, left_title, left_bullets),
            (6.7, right_title, right_bullets),
        ]:
            label = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.22), Inches(5.2), Inches(0.45))
            label.fill.solid()
            label.fill.fore_color.rgb = LIGHT
            label.line.color.rgb = RGBColor(215, 223, 230)
            tf = label.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            r = tf.paragraphs[0].add_run()
            r.text = section_title
            set_run_style(r, size=14, bold=True, color=NAVY)

            box = slide.shapes.add_textbox(Inches(x), Inches(1.85), Inches(5.2), Inches(4.95))
            t = box.text_frame
            t.word_wrap = True
            t.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            for idx, item in enumerate(bullets):
                level = 0
                text = item
                if isinstance(item, tuple):
                    level, text = item
                p = t.paragraphs[0] if idx == 0 else t.add_paragraph()
                p.level = level
                p.bullet = True
                p.space_after = Pt(7)
                run = p.add_run()
                run.text = text
                set_run_style(run, size=18 if level == 0 else 15, color=DARK)

    def add_table_slide(self, title, columns, rows, subtitle=None, note=None):
        slide = self.blank_slide()
        self.add_header(slide, title, subtitle)
        table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(0.55), Inches(1.35), Inches(12.2), Inches(4.95))
        table = table_shape.table
        col_width = Inches(12.2 / len(columns))
        for idx, name in enumerate(columns):
            table.columns[idx].width = col_width
            cell = table.cell(0, idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.text = name
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                set_run_style(run, size=11, bold=True, color=WHITE)
                p.alignment = PP_ALIGN.CENTER

        for r_idx, row in enumerate(rows, start=1):
            for c_idx, value in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if r_idx % 2 == 1 else WHITE
                cell.text = value
                p = cell.text_frame.paragraphs[0]
                for run in p.runs:
                    set_run_style(run, size=11, color=DARK)
                p.alignment = PP_ALIGN.LEFT

        if note:
            note_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.45), Inches(11.6), Inches(0.36))
            p = note_box.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = note
            set_run_style(run, size=10, color=MID)

    def add_flow_slide(self, title, steps, subtitle=None, note_items=None):
        slide = self.blank_slide()
        self.add_header(slide, title, subtitle)

        start_x = 0.5
        y = 2.0
        width = 1.55
        height = 0.9
        gap = 0.2
        for idx, step in enumerate(steps):
            x = start_x + idx * (width + gap)
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT
            box.line.color.rgb = BLUE
            tf = box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = step
            set_run_style(run, size=13, bold=True, color=NAVY)

            if idx < len(steps) - 1:
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x + width),
                    Inches(y + height / 2),
                    Inches(x + width + gap),
                    Inches(y + height / 2),
                )
                line.line.color.rgb = BLUE
                line.line.width = Pt(2)

        if note_items:
            note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(4.05), Inches(11.0), Inches(1.55))
            note.fill.solid()
            note.fill.fore_color.rgb = RGBColor(235, 244, 255)
            note.line.color.rgb = RGBColor(198, 219, 241)
            tf = note.text_frame
            tf.word_wrap = True
            for idx, item in enumerate(note_items):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.bullet = True
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = item
                set_run_style(run, size=16, color=DARK)

    def add_architecture_slide(self, title):
        slide = self.blank_slide()
        self.add_header(slide, title, "Core stack chosen in the architecture direction and already reflected in the repository")

        boxes = [
            (0.7, 1.55, 3.8, 0.85, BLUE, "Client layer", "Next.js website today; Android/iOS planned later"),
            (4.8, 1.55, 3.8, 0.85, TEAL, "API layer", "REST + OpenAPI 3.x contracts"),
            (8.9, 1.55, 3.8, 0.85, NAVY, "Runtime layer", "Spring Boot modular monolith"),
            (0.7, 3.0, 3.8, 1.1, LIGHT, "Data layer", "PostgreSQL 17 + Flyway migrations + JSON/JSONB-friendly config model"),
            (4.8, 3.0, 3.8, 1.1, LIGHT, "Security layer", "Signed bearer tokens, RBAC, admin MFA, audit trail"),
            (8.9, 3.0, 3.8, 1.1, LIGHT, "Delivery layer", "OCI images, Kubernetes overlays, OpenTofu/Terraform-compatible IaC, AWS workflow"),
            (0.7, 4.75, 12.0, 1.2, RGBColor(239, 246, 255), "Architectural principle", "Modular monolith first, microservices later only when justified; country-specific behavior must be modeled through generic configuration/data, never framework forks"),
        ]

        for x, y, w, h, color, title_text, body_text in boxes:
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            box.fill.solid()
            box.fill.fore_color.rgb = color
            box.line.color.rgb = RGBColor(205, 214, 223)
            tf = box.text_frame
            tf.word_wrap = True
            p1 = tf.paragraphs[0]
            r1 = p1.add_run()
            r1.text = title_text
            set_run_style(r1, size=15, bold=True, color=WHITE if color in (BLUE, TEAL, NAVY) else NAVY)
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = body_text
            set_run_style(r2, size=12, color=WHITE if color in (BLUE, TEAL, NAVY) else DARK)

    def add_repo_map_slide(self, title):
        slide = self.blank_slide()
        self.add_header(slide, title, "How the repository is organized today")

        tree = [
            "education-system-framework/",
            "├── backend/   Spring Boot parent with common, identity-access, organization, school-pack, attendance, gradebook, meal-catering, platform-core",
            "├── frontend/  independent frontend area with website active now and android/ios reserved",
            "├── devops/    container, kubernetes, IaC, automation scripts, AWS deployment assets",
            "├── df/        Dark Factory workflow, backlog, role docs, runtime tracking, task artifacts",
            "├── docs/      runbook, deployment guide, generated presentation assets",
            "├── design/    low-fidelity page artifacts for website tasks",
            "└── data/      source-backed country data packages such as the normalized Poland institution dataset",
        ]
        body = slide.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(12.0), Inches(3.25))
        tf = body.text_frame
        tf.word_wrap = True
        for idx, line in enumerate(tree):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.space_after = Pt(7)
            r = p.add_run()
            r.text = line
            set_run_style(r, size=16, color=DARK)

        note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.1), Inches(11.9), Inches(1.1))
        note.fill.solid()
        note.fill.fore_color.rgb = LIGHT
        note.line.color.rgb = RGBColor(214, 222, 230)
        tf2 = note.text_frame
        tf2.word_wrap = True
        p = tf2.paragraphs[0]
        r = p.add_run()
        r.text = "Key takeaway: the repository is not only application code. It also contains the operating framework, design assets, data packages, automation, and evidence needed to run an AI-governed SDLC."
        set_run_style(r, size=15, color=NAVY)

    def add_closing_slide(self, title, bullets, subtitle, closing_text):
        slide = self.blank_slide()
        self.add_header(slide, title, subtitle)
        body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(12.0), Inches(4.9))
        body.fill.solid()
        body.fill.fore_color.rgb = LIGHT
        body.line.color.rgb = RGBColor(214, 222, 230)
        tf = body.text_frame
        tf.word_wrap = True
        for idx, item in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.bullet = True
            p.space_after = Pt(10)
            r = p.add_run()
            r.text = item
            set_run_style(r, size=19, color=DARK)

        closing = slide.shapes.add_textbox(Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.3))
        p = closing.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = closing_text
        set_run_style(r, size=11, color=MID)


def build_detailed_deck():
    deck = DeckBuilder(DETAIL_FOOTER)

    deck.add_title_slide(
        "Education System Framework",
        [
            "Detailed project overview for stakeholders, contributors, and reviewers.",
            "Repository snapshot: 2026-05-27 • Dark Factory delivery model + sovereign education platform vision.",
        ],
        "Project, architecture, delivery model, and status",
    )

    deck.add_bullet_slide(
        "Executive summary",
        [
            "The project aims to build a sovereign, configurable, API-first education platform that countries can run on their own infrastructure.",
            "It combines two layers: the product itself (education-system framework) and the operating system for delivery (Dark Factory).",
            "Phase 1 foundations are largely in place: backend runtime, PostgreSQL/Flyway, OpenAPI, tenancy, configuration inheritance, audit, auth, RBAC, MFA, OCI/Kubernetes/AWS deployment paths, and an initial website shell.",
            "The next value step is turning these foundations into Phase 2 school-domain flows, starting with institution directory persistence/filtering and the broader Poland MVP feature set.",
        ],
    )

    deck.add_two_column_slide(
        "Problem and vision",
        "Why this project exists",
        [
            "Many education platforms are either proprietary SaaS, hard to customize safely, or too expensive to rebuild country by country.",
            "Countries need strong data residency and operational sovereignty, not dependence on a vendor-controlled production environment.",
            "Education structures vary significantly across countries, so one hard-coded domain model does not scale globally.",
        ],
        "What this framework proposes",
        [
            "A shared, generic core framework with country-specific differences expressed through configuration and data packages.",
            "A headless API-first architecture so website, mobile, desktop, AI assistants, and third parties can all consume the same contracts.",
            "Versioned release packages, migration guidance, compatibility checking, and country-operated deployment pipelines.",
        ],
    )

    deck.add_two_column_slide(
        "Primary stakeholders and users",
        "Administrative stakeholders",
        [
            "Ministry / country admin — governs country-wide rules, calendars, grade scales, and deployment-local policy.",
            "Regional / city admin — manages institutions in geographic scope without breaking country-level rules.",
            "Institution admin — operates the school/kindergarten/university configuration and local users.",
        ],
        "Operational end users",
        [
            "Teacher / lecturer — schedule, attendance, gradebook, homework, and later AI-assisted lesson preparation.",
            "Student — grades, homework, schedule, and guided AI assistance.",
            "Parent / guardian — child progress visibility, meal exclusions, communication, and oversight workflows.",
        ],
    )

    deck.add_bullet_slide(
        "Core principles and differentiators",
        [
            "Configuration over customization — the platform should adapt through settings, templates, and metadata rather than per-country code forks.",
            "Country-agnostic framework — country templates are data-only; no country may alter framework code, schema, or API contracts directly.",
            "API-first and headless — backend capabilities are exposed through documented contracts before client-specific polish.",
            "Sovereignty-aware deployment — each country owns its own infrastructure, data, backups, access, and deployment execution.",
            "Auditable and secure by default — auth, RBAC, MFA, immutable audit, and explicit runtime evidence are first-class design constraints.",
        ],
    )

    deck.add_flow_slide(
        "Dark Factory delivery model",
        ["Human request", "SA", "Design / lane", "QA", "PO", "Accepted output"],
        "The repository also documents how autonomous agents are expected to deliver work",
        [
            "Roles: SA, designer, backend-dev, frontend-dev, devops, data-engineer, QA, PO.",
            "Single-role-per-session keeps traceability and prevents self-approval.",
            "Work is only done after QA passes and PO accepts it.",
        ],
    )

    deck.add_table_slide(
        "Roadmap and delivery phases",
        ["Phase", "Goal", "Representative scope", "Status perspective"],
        [
            ["Phase 0", "Discovery", "Vision, backlog, architecture direction, Poland research", "Completed"],
            ["Phase 1", "Platform foundation", "Backend core, security, tenancy, config, audit, website shell, deployment baseline", "Largely implemented"],
            ["Phase 2", "Poland school MVP", "Institution hierarchy, subjects, schedule, attendance, grades, homework, meals, dashboards", "Partially prepared"],
            ["Phase 3", "Additional schema packs", "Kindergarten and university packs", "Planned"],
            ["Phase 4", "AI assistance", "Student and teacher AI assistants with governance rules", "Planned"],
            ["Phase 5-6", "Analytics and tracker integration", "Dashboards, statistics, issue-tracker sync", "Planned"],
        ],
        note="The current runtime board shows the strongest progress in Phase 1 foundations, with Phase 2 work beginning through design and data preparation tasks.",
    )

    deck.add_architecture_slide("Architecture and technology stack")
    deck.add_repo_map_slide("Repository and module map")

    deck.add_two_column_slide(
        "Implemented platform capabilities",
        "Backend and platform foundations",
        [
            "OpenAPI + Swagger routes, PostgreSQL/Flyway baseline, deployment-tenant model, configuration inheritance engine, compatibility/inheritance-break reporting, generic audit trail.",
            "Authentication foundation, role-based access control, and admin MFA are already implemented in backend modules.",
            "Translation storage is data-driven and database-backed, aligned with the no-language-specific-code rule.",
        ],
        "Delivery, data, and product-facing foundations",
        [
            "OCI image baseline, Kubernetes/OpenTofu-compatible deployment assets, and an on-demand AWS GitHub Actions workflow are in the repository.",
            "The website frontend exists as an independent Next.js + React project with home/login/student/teacher routes and low-fidelity UI assets.",
            "A normalized Poland institution dataset has been prepared for future selector/filtering and persistence work.",
        ],
    )

    deck.add_two_column_slide(
        "Runnable system and operator experience today",
        "What can run locally",
        [
            "`compose.local.yaml` starts PostgreSQL, the Spring backend, and the `frontend/website` app together.",
            "Backend entrypoint: `backend/platform-core`; main health route: `GET /platform/status`.",
            "API docs are exposed via `GET /api-docs` and `GET /swagger-ui`.",
            "Website routes currently include `/`, `/login`, `/student`, and `/teacher`.",
        ],
        "Deployment and operations",
        [
            "DevOps assets support OCI image builds and a portable Kubernetes/IaC baseline for AWS, Azure, Google Cloud, and self-hosted targets.",
            "The first live cloud path is a manual GitHub Actions workflow that builds `platform-core`, pushes to ECR, and deploys to EKS.",
            "Launcher scripts automate Dark Factory role handoffs for JetBrains/Git Bash/PowerShell operator workflows.",
        ],
    )

    deck.add_bullet_slide(
        "Governance and non-negotiable guardrails",
        [
            "Country templates are data-only. If a country requirement cannot be modeled through generic configuration/data, it must trigger architecture review rather than a one-off code fork.",
            "Real place, school, and subject names must be source-backed; teacher names, student names, and grade records must remain synthetic in country/test datasets.",
            "UI-facing frontend implementation requires a designer package first; frontend must not invent UI from scratch when design evidence is missing.",
            "One session = one role. This keeps design, implementation, QA, and PO acceptance separated and auditable.",
            "Nothing is considered done until QA passes and PO explicitly accepts the result.",
        ],
    )

    deck.add_table_slide(
        "Runtime status snapshot on 2026-05-27",
        ["Bucket", "Items", "Meaning"],
        [
            ["Done foundation work", "STORY-010/011/012/013/014/020/021/022/023/030/031/040/050/080/081/082/220 and TASK-007", "Core platform, security, deployment, documentation, and workflow foundations are accepted"],
            ["Ready for QA", "TASK-008, TASK-009, TASK-010, TASK-012, TASK-013", "Documentation, local startup, AWS deployment workflow, launcher automation, and presentation assets need independent verification"],
            ["Ready for PO", "TASK-011", "The Poland institution dataset package passed QA and awaits product acceptance"],
            ["Blocked", "TASK-006", "Website validation is blocked on a workstation missing Node.js/npm tooling"],
        ],
        note="This slide intentionally separates accepted work from work still awaiting QA/PO so the presentation does not overstate completion.",
    )

    deck.add_bullet_slide(
        "Key risks and constraints",
        [
            "Workstation tooling can still block delivery; the active example is missing Node.js/npm for `TASK-006` validation.",
            "Phase 2 scope is large. The project must continue protecting Phase 1 foundations from scope creep into unfinished school-domain breadth too early.",
            "Some architecture decisions are intentionally deferred, such as schema-isolation strategy, distributed caching needs, AI provider selection, and deeper infrastructure choices.",
            "Deployment remains country-operated by design, so production readiness depends on operator-managed secrets, observability, databases, and cloud account access outside the repository.",
        ],
    )

    deck.add_closing_slide(
        "Recommended next steps",
        [
            "Close the current runtime queue by finishing QA and PO reviews for tasks 008 through 013.",
            "Unblock `TASK-006` on a workstation with Node.js 20+ so the website low-fidelity implementation can complete lint, type-check, build, and manual login-flow validation.",
            "Move the accepted Poland institution dataset into an end-to-end persistence and homepage-filter flow through backend and website tasks.",
            "Use the accepted Phase 1 platform foundation to start the highest-value school-domain stories in Phase 2: institution hierarchy, people management, subjects, attendance, gradebook, and dashboards.",
        ],
        "Recommended sequencing after this presentation",
        "This deck is designed to be updated as the runtime board moves from platform foundation into the school-domain MVP.",
    )

    return deck.prs


def build_executive_ru_deck():
    deck = DeckBuilder(EXECUTIVE_RU_FOOTER)

    deck.add_title_slide(
        "Education Framework System",
        [
            "Короткая executive-презентация для коллег и руководства.",
            "Фокус: бизнес-боли, ценность продукта и общее техническое описание платформы.",
        ],
        "Единая цифровая платформа для системы образования",
    )

    deck.add_bullet_slide(
        "Общее описание проекта",
        [
            "Единая цифровая платформа для министерства, школ, садов, университетов.",
            "Объединяет управление, обучение, коммуникацию и цифровую инфраструктуру.",
            "Один framework вместо набора разрозненных локальных решений.",
            "Основа для масштабирования на разные страны через конфигурацию и данные.",
        ],
        "Не отдельное приложение, а платформенная основа для всей системы образования",
    )

    deck.add_bullet_slide(
        "Почему проект нужен сейчас",
        [
            "Фрагментированные данные и процессы.",
            "Нет real-time картины по системе образования.",
            "Разный уровень цифровизации учреждений.",
            "Высокая операционная нагрузка на школы, учителей и родителей.",
            "Слишком много ручной координации и локальных решений.",
        ],
        "Сегодня цифровая среда образования неоднородна и плохо управляется",
    )

    deck.add_two_column_slide(
        "Боли государства и учреждений",
        "Министерство / органы управления",
        [
            "Нет единого окна управления учреждениями.",
            "Нет полной статистики в реальном времени.",
            "Сложно централизованно запускать коммуникации и инициативы.",
            "Трудно держать единый цифровой стандарт по стране.",
        ],
        "Школы / сады / университеты",
        [
            "Разрозненные сайты или их отсутствие.",
            "Локальная поддержка контента и документов за свой счет.",
            "Слабая онлайн-видимость событий, кружков, объявлений.",
            "Лишние расходы на несвязанные между собой решения.",
        ],
        "Главная проблема: низкая управляемость и высокая стоимость разрозненности",
    )

    deck.add_bullet_slide(
        "Боли учителей, родителей и учеников",
        [
            "Учитель: перегрузка рутиной, подготовкой, оценками, коммуникацией.",
            "Родитель: вынужден объяснять темы и компенсировать пробелы системы.",
            "Ученик: нет одной понятной точки входа по задачам, расписанию и рискам.",
            "Проблемы замечаются поздно, вовлеченность падает.",
            "Нагрузка переносится с системы на семью и учителя.",
        ],
    )

    deck.add_bullet_slide(
        "Наш ответ: единая платформа",
        [
            "Dashboards для министерства, школы, учителя, родителя, ученика.",
            "Единое управление учреждениями и структурой системы образования.",
            "Школьные сайты, документы, события, кружки, объявления в одной среде.",
            "AI-assistant для ученика и учителя.",
            "Цифровые учебники и единый доступ к контенту.",
        ],
        "Один продукт закрывает управленческую, учебную и коммуникационную части",
    )

    deck.add_bullet_slide(
        "Общее техническое описание проекта",
        [
            "Backend: Java + Spring Boot, модульная архитектура.",
            "Frontend: Next.js / React для web; mobile можно добавлять позже.",
            "База данных: PostgreSQL + Flyway migrations.",
            "API-first подход: OpenAPI / Swagger.",
            "Инфраструктура: OCI containers, Kubernetes, IaC/OpenTofu, AWS deployment path.",
            "Конфигурация и данные вместо country-specific кода.",
        ],
        "Технически это платформенный framework, а не одноразовый проект под одного клиента",
    )

    deck.add_two_column_slide(
        "Ключевая ценность",
        "Для системы образования",
        [
            "Прозрачность в реальном времени.",
            "Централизованное управление.",
            "Снижение стоимости владения.",
            "Единый стандарт качества.",
        ],
        "Для пользователей",
        [
            "Учитель — меньше рутины.",
            "Родитель — больше прозрачности.",
            "Ученик — понятная среда и помощь в обучении.",
            "Школа — готовая цифровая инфраструктура.",
        ],
        "Ценность измеряется снижением хаоса, затрат и ручной нагрузки",
    )

    deck.add_bullet_slide(
        "Дополнительная боль: нестабильность аутсорса",
        [
            "Аутсорс = зависимость от внешних заказов.",
            "Нет уверенности в завтрашнем дне у команды и бизнеса.",
            "Чужие приоритеты определяют судьбу компании.",
            "Собственный продукт = стратегический актив.",
            "Если продукт решает универсальную боль, он масштабируется глобально.",
        ],
        "Переход от продажи часов к созданию собственного международного продукта",
    )

    deck.add_bullet_slide(
        "Почему framework-подход сильнее",
        [
            "Одна архитектура для многих стран и учреждений.",
            "Гибкость через конфигурацию и данные.",
            "Быстрее масштабирование и ниже стоимость изменений.",
            "Меньше одноразовых кастомизаций и форков.",
            "Предсказуемое развитие продукта.",
        ],
    )

    deck.add_two_column_slide(
        "Внедрение и масштабирование",
        "Как запускаем",
        [
            "Предзагрузка структуры регионов и учреждений.",
            "Демо-среда с синтетическими данными.",
            "Импорт реальных учеников и учителей на этапе внедрения.",
            "Гайды, туториалы, видеообучение.",
        ],
        "Как масштабируем",
        [
            "Новые учреждения без новой архитектуры.",
            "Новые страны без country-specific форков.",
            "Единая платформа для дальнейших AI и analytics сценариев.",
            "Рост через повторяемую модель, а не через проектную разработку с нуля.",
        ],
        "Снижаем порог входа и ускоряем масштабирование",
    )

    deck.add_closing_slide(
        "Итог",
        [
            "Единая цифровая операционная система для образования.",
            "Решает боли государства, учреждений, учителей, родителей и учеников.",
            "Даёт бизнесу переход от нестабильного аутсорса к собственному продукту.",
            "Имеет потенциал масштабирования на международный рынок.",
        ],
        "Короткая версия для руководства и коллег",
        "Ключевая идея: продукт одновременно решает системную проблему образования и создает устойчивую платформенную бизнес-модель.",
    )

    return deck.prs


def save_deck(build_fn, output_name):
    base_dir = Path(__file__).resolve().parent
    output_path = base_dir / output_name
    presentation = build_fn()
    presentation.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Slides: {len(presentation.slides)}")
    return output_path, len(presentation.slides)


if __name__ == "__main__":
    save_deck(build_detailed_deck, DETAIL_OUTPUT_NAME)
    save_deck(build_executive_ru_deck, EXECUTIVE_RU_OUTPUT_NAME)

